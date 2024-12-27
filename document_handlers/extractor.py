import logging
from fastapi import HTTPException, UploadFile
from PyPDF2 import PdfReader
import tiktoken
from io import BytesIO
import pandas as pd
from typing import Tuple, List
from docx import Document
from pptx import Presentation
import zipfile
import os
import tempfile
import aiofiles
from pathlib import Path

logger = logging.getLogger(__name__)

class DocumentExtractor:
    SUPPORTED_TYPES = {
        'application/pdf': 'pdf',
        'application/vnd.ms-excel': 'excel',  # .xls
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'excel',  # .xlsx
        'text/plain': 'text',
        'text/csv': 'csv',
        'application/msword': 'word',  # .doc
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'word',  # .docx
        'application/vnd.ms-powerpoint': 'powerpoint',  # .ppt
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'powerpoint',  # .pptx
        'application/zip': 'zip',
        'application/x-zip-compressed': 'zip'
    }

    SUPPORTED_EXTENSIONS = {
        '.txt', '.csv', '.pdf', '.xlsx', '.xls', 
        '.doc', '.docx', '.ppt', '.pptx'
    }

    def __init__(self):
        self.encoding = tiktoken.get_encoding("cl100k_base")
        self.max_tokens = 6000

    def _get_file_type(self, file: UploadFile) -> str:
        """Determine the file type from the uploaded file."""
        content_type = file.content_type
        file_type = self.SUPPORTED_TYPES.get(content_type)
        
        if not file_type:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {content_type}. Supported types: {list(self.SUPPORTED_TYPES.keys())}"
            )
        
        return file_type

    async def extract_text(self, file: UploadFile) -> Tuple[str, str]:
        """
        Extract text from the uploaded file based on its type.
        Returns: Tuple[extracted_text, file_type]
        """
        try:
            file_type = self._get_file_type(file)
            logger.info(f"Processing {file_type} file: {file.filename}")
            
            content = await file.read()
            
            # Extract text based on file type
            if file_type == 'zip':
                text = await self._extract_zip(content, file.filename)
            elif file_type == 'pdf':
                text = await self._extract_pdf(content, file.filename)
            elif file_type in ['excel', 'csv']:
                text = await self._extract_tabular(content, file.filename, file_type)
            elif file_type == 'word':
                text = await self._extract_word(content, file.filename)
            elif file_type == 'powerpoint':
                text = await self._extract_powerpoint(content, file.filename)
            elif file_type == 'text':
                text = content.decode('utf-8')
            
            # Ensure token limit
            if len(self.encoding.encode(text)) > self.max_tokens:
                truncated_tokens = self.encoding.encode(text)[:self.max_tokens]
                text = self.encoding.decode(truncated_tokens)
            
            logger.info(f"Successfully extracted text from {file_type} file: {file.filename} "
                       f"(tokens: {len(self.encoding.encode(text))})")
            return text, file_type
            
        except Exception as e:
            logger.error(f"Error extracting text from {file.filename}: {str(e)}", exc_info=True)
            raise HTTPException(status_code=500, detail=str(e))

    async def _extract_pdf(self, content: bytes, filename: str) -> str:
        """Extract text from PDF content."""
        pdf = PdfReader(BytesIO(content))
        text = ""
        
        for page in pdf.pages:
            page_text = page.extract_text()
            current_text = text + page_text
            
            if len(self.encoding.encode(current_text)) > self.max_tokens:
                remaining_tokens = self.max_tokens - len(self.encoding.encode(text))
                if remaining_tokens > 0:
                    truncated_tokens = self.encoding.encode(page_text)[:remaining_tokens]
                    text += self.encoding.decode(truncated_tokens)
                break
            else:
                text += page_text
                
        return text

    async def _extract_tabular(self, content: bytes, filename: str, file_type: str) -> str:
        """Extract text from Excel or CSV files."""
        try:
            if file_type == 'excel':
                # Try openpyxl engine first (for .xlsx)
                try:
                    df = pd.read_excel(BytesIO(content), engine='openpyxl')
                except Exception as e:
                    # If that fails, try xlrd engine (for .xls)
                    logger.info(f"Falling back to xlrd engine for {filename}: {str(e)}")
                    df = pd.read_excel(BytesIO(content), engine='xlrd')
            else:  # csv
                df = pd.read_csv(BytesIO(content))
            
            # Convert DataFrame to a more readable format
            return df.to_string(index=False)  # Exclude index for cleaner output
            
        except Exception as e:
            logger.error(f"Error processing tabular file {filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error processing {file_type} file: {str(e)}"
            )

    async def _extract_word(self, content: bytes, filename: str) -> str:
        """Extract text from Word document content."""
        try:
            doc = Document(BytesIO(content))
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                paragraph_text = paragraph.text + "\n"
                current_text = text + paragraph_text
                
                # Check token limit
                if len(self.encoding.encode(current_text)) > self.max_tokens:
                    remaining_tokens = self.max_tokens - len(self.encoding.encode(text))
                    if remaining_tokens > 0:
                        truncated_tokens = self.encoding.encode(paragraph_text)[:remaining_tokens]
                        text += self.encoding.decode(truncated_tokens)
                    break
                else:
                    text += paragraph_text
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells) + "\n"
                    current_text = text + row_text
                    
                    if len(self.encoding.encode(current_text)) > self.max_tokens:
                        break
                    text += row_text
                
                if len(self.encoding.encode(text)) > self.max_tokens:
                    break
            
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error processing Word document {filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error processing Word document: {str(e)}"
            )

    async def _extract_zip(self, content: bytes, filename: str) -> str:
        """Extract text from supported files within a ZIP archive."""
        try:
            # Create a temporary directory to extract files
            with tempfile.TemporaryDirectory() as temp_dir:
                # Save ZIP content to a temporary file
                zip_path = os.path.join(temp_dir, 'temp.zip')
                async with aiofiles.open(zip_path, 'wb') as f:
                    await f.write(content)

                # Extract all files
                with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)

                # Process each file in the ZIP
                all_text = []
                total_tokens = 0
                
                for root, _, files in os.walk(temp_dir):
                    for file in files:
                        if file == 'temp.zip':
                            continue
                            
                        file_path = os.path.join(root, file)
                        extension = Path(file).suffix.lower()
                        
                        if extension not in self.SUPPORTED_EXTENSIONS:
                            logger.info(f"Skipping unsupported file: {file}")
                            continue

                        try:
                            # Read file content
                            async with aiofiles.open(file_path, 'rb') as f:
                                file_content = await f.read()

                            # Extract text based on file extension
                            if extension in ['.pdf']:
                                file_text = await self._extract_pdf(file_content, file)
                            elif extension in ['.xlsx', '.xls', '.csv']:
                                file_text = await self._extract_tabular(
                                    file_content, 
                                    file,
                                    'excel' if extension in ['.xlsx', '.xls'] else 'csv'
                                )
                            elif extension in ['.doc', '.docx']:
                                file_text = await self._extract_word(file_content, file)
                            else:  # .txt files
                                file_text = file_content.decode('utf-8')

                            # Add file header
                            file_text = f"\n=== {file} ===\n{file_text}"
                            
                            # Check token limit
                            new_tokens = len(self.encoding.encode(file_text))
                            if total_tokens + new_tokens > self.max_tokens:
                                logger.warning(f"Token limit reached. Skipping remaining files in ZIP.")
                                break
                            
                            total_tokens += new_tokens
                            all_text.append(file_text)

                        except Exception as e:
                            logger.error(f"Error processing file {file} in ZIP: {str(e)}")
                            all_text.append(f"\n=== Error processing {file}: {str(e)} ===\n")

                return "\n".join(all_text)

        except Exception as e:
            logger.error(f"Error processing ZIP file {filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error processing ZIP file: {str(e)}"
            )

    async def _extract_powerpoint(self, content: bytes, filename: str) -> str:
        """Extract text from PowerPoint presentation content."""
        try:
            prs = Presentation(BytesIO(content))
            text = []
            
            # Process each slide
            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n=== Slide {slide_number} ===\n"]
                
                # Extract text from shapes (including title and content placeholders)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_text = shape.text.strip()
                        
                        # Check token limit
                        current_text = "\n".join(text) + "\n".join(slide_text) + shape_text
                        if len(self.encoding.encode(current_text)) > self.max_tokens:
                            logger.warning(f"Token limit reached at slide {slide_number}. Truncating...")
                            text.extend(slide_text)
                            return "\n".join(text)
                        
                        slide_text.append(shape_text)
                
                # Add notes if they exist
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_text = f"[Notes: {slide.notes_slide.notes_text_frame.text.strip()}]"
                    current_text = "\n".join(text) + "\n".join(slide_text) + notes_text
                    if len(self.encoding.encode(current_text)) <= self.max_tokens:
                        slide_text.append(notes_text)
                
                text.extend(slide_text)
            
            return "\n".join(text)
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error processing PowerPoint file: {str(e)}"
            )